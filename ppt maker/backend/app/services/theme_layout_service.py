"""Theme tokens and slide styling for python-pptx exports."""

from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt

# Structured deck flow: title → agenda → content slides → conclusion
DECK_FLOW = ("title", "agenda", "content", "conclusion")


class ThemeLayoutService:
    """Maps SaaS theme names to PowerPoint colors and typography."""

    @staticmethod
    def theme_colors(theme: str) -> tuple[RGBColor, RGBColor]:
        if theme.lower().startswith("cinematic"):
            return RGBColor(15, 23, 42), RGBColor(241, 245, 249)
        return RGBColor(255, 255, 255), RGBColor(30, 41, 59)

    def style_slide_background(self, slide, theme: str) -> None:
        bg, _ = self.theme_colors(theme)
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = bg

    def style_title(self, title_shape, theme: str) -> None:
        _, text_color = self.theme_colors(theme)
        paragraph = title_shape.text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.font.bold = True
        paragraph.font.size = Pt(36)
        paragraph.font.color.rgb = text_color

    def style_subtitle(self, subtitle_shape, theme: str) -> None:
        _, text_color = self.theme_colors(theme)
        paragraph = subtitle_shape.text_frame.paragraphs[0]
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = text_color

    def style_body(self, text_frame, theme: str) -> None:
        _, text_color = self.theme_colors(theme)
        for paragraph in text_frame.paragraphs:
            paragraph.font.color.rgb = text_color
