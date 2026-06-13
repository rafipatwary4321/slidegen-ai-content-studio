from pathlib import Path
from uuid import uuid4

from pptx import Presentation
from pptx.util import Pt

from app.models.schemas import ExportResult, SlideOutlineItem, SpeakerNoteItem
from app.services.theme_layout_service import ThemeLayoutService


class PptxService:
    """Generates native PowerPoint files: title → agenda → content → conclusion."""

    def __init__(self) -> None:
        self.output_dir = Path(__file__).resolve().parents[2] / "output" / "pptx"
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self._theme = ThemeLayoutService()

    def export(
        self,
        title: str,
        slides: list[SlideOutlineItem],
        theme: str,
        speaker_notes: list[SpeakerNoteItem] | None = None,
    ) -> ExportResult:
        export_id = str(uuid4())
        safe_title = self._safe_file_name(title)
        file_path = self.output_dir / f"{export_id}_{safe_title}.pptx"

        presentation = Presentation()
        self._add_title_slide(presentation, title, theme)
        self._add_agenda_slide(presentation, slides, theme)
        self._add_content_slides(presentation, slides, speaker_notes or [], theme)
        self._add_conclusion_slide(presentation, slides, theme)
        presentation.save(file_path)

        return ExportResult(
            export_id=export_id,
            title=title,
            theme=theme,
            slide_count=len(presentation.slides),
            file_path=str(file_path),
            download_url=f"/api/v1/exports/pptx/{export_id}/download",
            status="ready",
        )

    def resolve_export_path(self, export_id: str) -> Path | None:
        matches = sorted(self.output_dir.glob(f"{export_id}_*.pptx"))
        return matches[0] if matches else None

    def _add_title_slide(self, presentation: Presentation, title: str, theme: str) -> None:
        slide = presentation.slides.add_slide(presentation.slide_layouts[0])
        title_shape = slide.shapes.title
        subtitle = slide.placeholders[1]
        title_shape.text = title
        subtitle.text = "SlideGen AI"
        self._theme.style_slide_background(slide, theme)
        self._theme.style_title(title_shape, theme)
        self._theme.style_subtitle(subtitle, theme)

    def _add_agenda_slide(self, presentation: Presentation, slides: list[SlideOutlineItem], theme: str) -> None:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "Agenda"
        body = slide.placeholders[1].text_frame
        body.clear()
        for idx, item in enumerate(slides[:8], start=1):
            p = body.paragraphs[0] if idx == 1 else body.add_paragraph()
            p.text = f"{idx}. {item.title}"
            p.level = 0
            p.font.size = Pt(22)
        self._theme.style_slide_background(slide, theme)
        self._theme.style_title(slide.shapes.title, theme)
        self._theme.style_body(body, theme)

    def _add_content_slides(
        self,
        presentation: Presentation,
        slides: list[SlideOutlineItem],
        speaker_notes: list[SpeakerNoteItem],
        theme: str,
    ) -> None:
        notes_map = {n.slide_number: n.note for n in speaker_notes}
        for item in slides:
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes.title.text = item.title
            body = slide.placeholders[1].text_frame
            body.clear()

            bullets = item.bullets[:5]
            for idx, bullet in enumerate(bullets, start=1):
                text = bullet if len(bullet) <= 140 else f"{bullet[:137]}..."
                p = body.paragraphs[0] if idx == 1 else body.add_paragraph()
                p.text = text
                p.level = 0
                p.font.size = Pt(20)
                p.space_after = Pt(8)

            note_text = notes_map.get(item.slide_number)
            if note_text:
                slide.notes_slide.notes_text_frame.text = note_text

            self._theme.style_slide_background(slide, theme)
            self._theme.style_title(slide.shapes.title, theme)
            self._theme.style_body(body, theme)

    def _add_conclusion_slide(self, presentation: Presentation, slides: list[SlideOutlineItem], theme: str) -> None:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "Conclusion & Next Steps"
        body = slide.placeholders[1].text_frame
        body.clear()
        highlights = [s.title for s in slides[:3]] or ["Review outline", "Refine with prompts", "Export deck"]
        for idx, text in enumerate(highlights, start=1):
            p = body.paragraphs[0] if idx == 1 else body.add_paragraph()
            p.text = f"Takeaway: {text}"
            p.level = 0
            p.font.size = Pt(22)
        self._theme.style_slide_background(slide, theme)
        self._theme.style_title(slide.shapes.title, theme)
        self._theme.style_body(body, theme)

    @staticmethod
    def _safe_file_name(title: str) -> str:
        cleaned = "".join(ch if ch.isalnum() or ch in {"-", "_"} else "_" for ch in title.strip().replace(" ", "_"))
        return cleaned[:60] or "presentation"
